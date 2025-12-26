import os
from pathlib import Path
import json
import io
import pathlib
import csv
from functools import partial
import logging
# 3rd Party
import requests

logger = logging.getLogger("Parsers")

# ALL VALID IMAGE EXTENSIONS:
# [".png", ".jpg", ".jpeg", ".gif", ".webp", ".heic", ".heif", ".tif", ".tiff", ".bmp", ".ico"]
# ALL VALID TEXT EXTENSIONS:
# [".txt", ".md", ".pdf", ".docx", ".gdoc", ".rtf", ".pptx", ".csv", ".xlsx", ".xls", ".json", ".yaml", ".yml", ".xml", ".ini", ".toml", ".env", ".py", ".js", ".jsx", ".ts", ".tsx", ".html", ".css", ".c", ".cpp", ".h", ".java", ".cs", ".php", ".rb", ".go", ".rs", ".sql", ".sh", ".bat", ".ps1"]

# --- OPTIONAL IMPORTS (Soft Dependencies) ---
try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pdfminer.high_level import extract_text
except ImportError:
    extract_text = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import pillow_heif
    pillow_heif.register_heif_opener()
except ImportError:
    pass

BASE_DIR = Path(os.path.dirname(os.path.abspath(__file__)))
DATA_DIR = Path(os.getenv('LOCALAPPDATA')) / "2nd Brain"

# --- HELPERS ---

def is_connected():
    try:
        requests.head('http://www.google.com', timeout=1)
        return True
    except (requests.ConnectionError, requests.Timeout):
        return False

# --- DRIVE AUTHENTICATION ---

def get_drive_service(config):
    """
    Handles OAuth 2.0 flow. 
    Returns the Google Drive service object or None if disabled/offline.
    """
    # 1. Check Config
    if not config.get("use_drive", False):
        logger.info("[Drive] Disabled by config.")
        return None
    
    # 2. Check Internet
    if not is_connected():
        logger.warning("[Drive] No internet — skipping.")
        return None
    
    # 3. Import Google Libs (Lazy Import)
    try:
        from google.auth.transport.requests import Request
        from google.oauth2.credentials import Credentials
        from google_auth_oauthlib.flow import InstalledAppFlow
        from googleapiclient.discovery import build
    except ImportError:
        logger.error("[Drive] Google client libraries not installed.")
        return None

    # 4. Check Credentials.json
    cred_path = DATA_DIR / "credentials.json"
    # This is a good place to put it.
        
    if not cred_path.exists():
        logger.error("[Drive] No credentials.json found.")
        return None

    SCOPES = ["https://www.googleapis.com/auth/drive.readonly"]
    creds = None
    token_path = DATA_DIR / "token.json"
    # This is mutable data, so it should go in DATA_DIR.

    # 5. Load Token
    if token_path.exists():
        try:
            creds = Credentials.from_authorized_user_file(str(token_path), SCOPES)
        except Exception:
            pass # Invalid token, will refresh/re-auth

    # 6. Refresh or Login
    try:
        if not creds or not creds.valid:
            if creds and creds.expired and creds.refresh_token:
                creds.refresh(Request())
            else:
                logger.info("[Drive] Authenticating in browser...")
                flow = InstalledAppFlow.from_client_secrets_file(str(cred_path), SCOPES)
                creds = flow.run_local_server(port=0)
            
            # Save token
            with open(token_path, "w") as token:
                token.write(creds.to_json())
        
        service = build("drive", "v3", credentials=creds, cache_discovery=False)
        logger.info("[Drive] Authenticated successfully.")
        return service
        
    except Exception as e:
        logger.error(f"[Drive] Auth failed: {e}")
        return None

def download_drive_content(drive_service, doc_id: str, mimeType: str):
    """Downloads a Google Doc's content as plain text using its file ID."""
    try:
        from googleapiclient.http import MediaIoBaseDownload
        request = drive_service.files().export_media(fileId=doc_id, mimeType=mimeType)
        
        fh = io.BytesIO()
        downloader = MediaIoBaseDownload(fh, request)
        
        done = False
        while not done:
            status, done = downloader.next_chunk()
        
        fh.seek(0)
        text = fh.read().decode('utf-8')
        logger.info(f"[Drive Download] {doc_id}: {len(text)} bytes")
        return text
    except Exception as e:
        logger.error(f"[Drive Download Error] {e}")
        return None

# --- FILE PARSERS ---

def parse_gdoc(file_path: pathlib.Path, drive_service, max_chars: int) -> str:
    """Parses a .gdoc file (JSON shortcut) and fetches content."""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            gdoc_data = json.load(f)
        
        doc_url = gdoc_data.get('doc_id')
        if not doc_url:
            return None # Invalid shortcut
            
        if not drive_service:
            return None # Cannot parse without service

        content = download_drive_content(drive_service, doc_url, "text/plain")
        
        # If content is None (download failed), return None so we don't mark as done
        if content is None:
            return None
            
        return content[:max_chars]
        
    except Exception as e:
        logger.error(f"[Parser Error] parse_gdoc {file_path.name}: {e}")
        return None

def parse_docx(file_path: pathlib.Path, max_chars: int) -> str:
    if not Document: return ""
    try:
        doc = Document(file_path)
        full_text = []
        current_len = 0
        for para in doc.paragraphs:
            text = para.text
            full_text.append(text)
            current_len += len(text)
            if current_len > max_chars:
                break
        return '\n'.join(full_text)[:max_chars]
    except Exception as e:
        logger.error(f"[Parser Error] parse_docx: {file_path.name}: {e}")
        return ""

def parse_pptx(file_path: pathlib.Path, max_chars: int) -> str:
    if not Presentation: return ""
    try:
        prs = Presentation(file_path)
        text_runs = []
        current_len = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text
                    text_runs.append(text)
                    current_len += len(text)
                    if current_len > max_chars:
                        return "\n".join(text_runs)[:max_chars]
        return "\n".join(text_runs)
    except Exception as e: 
        logger.error(f"[Parser Error] parse_pptx: {file_path.name}: {e}")
        return ""

def parse_pdf(file_path: pathlib.Path, max_chars: int) -> str:
    if not extract_text: return ""
    try:
        text = extract_text(file_path)
        return text[:max_chars]
    except Exception as e: 
        logger.error(f"[Parser Error] parse_pdf: {file_path.name}: {e}")
        return ""

def parse_code_or_text(file_path: pathlib.Path, max_chars: int) -> str:
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read(max_chars)
    except Exception:
        logger.error(f"[Parser Error] parse_code_or_text: {file_path.name}")
        return ""

def parse_csv(file_path: pathlib.Path, max_chars: int) -> str:
    try:
        text_output = []
        current_len = 0
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            reader = csv.reader(f)
            try:
                headers = next(reader)
            except StopIteration:
                return ""

            for row in reader:
                row_str = ", ".join([f"{h}: {v}" for h, v in zip(headers, row) if v.strip()])
                text_output.append(row_str)
                current_len += len(row_str)
                if current_len > max_chars: break
            
        return "\n".join(text_output)
    except Exception as e:
        logger.error(f"[Parser Error] parse_csv: {file_path.name}: {e}")
        return ""

def parse_xlsx(file_path: pathlib.Path, max_chars: int) -> str:
    if not openpyxl: return ""
    try:
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        sheet = wb.active
        text_output = []
        current_len = 0
        rows = sheet.iter_rows(values_only=True)
        try:
            headers = next(rows)
        except StopIteration: return ""
            
        for row in rows:
            row_str = ", ".join([f"{h}: {v}" for h, v in zip(headers, row) if h and v])
            if row_str:
                text_output.append(row_str)
                current_len += len(row_str)
            if current_len > max_chars: break
                
        return "\n".join(text_output)
    except Exception as e: 
        logger.error(f"[Parser Error] parse_xlsx: {file_path.name}: {e}")
        return ""

def parse_image_placeholder(file_path: pathlib.Path, max_chars: int = 0) -> str:
    return "[IMAGE]"

# --- GLOBAL MAPPING (Optimization) ---
_EXTENSION_MAPPING = {
    '.txt': parse_code_or_text,
    '.md': parse_code_or_text,
    '.markdown': parse_code_or_text,
    '.docx': parse_docx,
    '.pdf': parse_pdf,
    '.rtf': parse_code_or_text,
    '.pptx': parse_pptx,
    '.csv': parse_csv,
    '.xlsx': parse_xlsx,
    '.xls': parse_xlsx,
    '.json': parse_code_or_text,
    '.yaml': parse_code_or_text,
    '.yml': parse_code_or_text,
    '.xml': parse_code_or_text,
    '.ini': parse_code_or_text,
    '.toml': parse_code_or_text,
    '.env': parse_code_or_text,
    '.py': parse_code_or_text,
    '.js': parse_code_or_text,
    '.jsx': parse_code_or_text,
    '.ts': parse_code_or_text,
    '.tsx': parse_code_or_text,
    '.html': parse_code_or_text,
    '.css': parse_code_or_text,
    '.c': parse_code_or_text,
    '.cpp': parse_code_or_text,
    '.h': parse_code_or_text,
    '.java': parse_code_or_text,
    '.cs': parse_code_or_text,
    '.php': parse_code_or_text,
    '.rb': parse_code_or_text,
    '.go': parse_code_or_text,
    '.rs': parse_code_or_text,
    '.sql': parse_code_or_text,
    '.sh': parse_code_or_text,
    '.bat': parse_code_or_text,
    '.ps1': parse_code_or_text,
}

def file_handler(extension: str, is_multimodal: bool, use_drive: bool, config: dict):
    """
    Returns a partial function with max_chars bound.
    The returned function signature is effectively: func(path) -> str
    (Except for gdoc, which still expects drive_service args via special handling in the loop)
    """
    try:
        ext = extension.lower()
        
        # 1. Check Global Map
        func = _EXTENSION_MAPPING.get(ext, None)
        
        # 2. Overrides
        if is_multimodal and ext in config['image_extensions']:
            func = parse_image_placeholder
                
        if use_drive and ext == '.gdoc':
            func = parse_gdoc
        
        return func
    except Exception as e:
        logger.error(f"[Parser Error] file_handler: {e}")
        return None